"""
PPTX Builder for Operation Quality Report - Simplified Design.

Creates professional PowerPoint presentations with one message per slide,
featuring chart + AI insight layouts in a restrained 2B business style.
"""

import io
import base64
from datetime import datetime
from typing import Dict, List, Any, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config import OUTPUT_DIR
from i18n import get_all_texts


# =============================================================================
# Professional 2B Color Palette - Restrained & Elegant
# =============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


# Dark theme colors
BG_DARK = hex_to_rgb("#0f172a")       # Deep navy background
BG_CARD = hex_to_rgb("#1e293b")       # Card background
BG_ACCENT = hex_to_rgb("#334155")     # Accent areas

# Brand colors
PRIMARY = hex_to_rgb("#3b82f6")       # Professional blue
SUCCESS = hex_to_rgb("#22c55e")       # Green
WARNING = hex_to_rgb("#eab308")       # Amber
DANGER = hex_to_rgb("#ef4444")        # Red

# Text colors
TEXT_WHITE = hex_to_rgb("#f8fafc")    # Primary text
TEXT_GRAY = hex_to_rgb("#94a3b8")     # Secondary text
TEXT_MUTED = hex_to_rgb("#64748b")    # Muted text

# Process colors
COLOR_INCIDENT = hex_to_rgb("#3b82f6")   # Blue
COLOR_CHANGE = hex_to_rgb("#22c55e")     # Green
COLOR_REQUEST = hex_to_rgb("#f59e0b")    # Amber
COLOR_PROBLEM = hex_to_rgb("#ec4899")    # Pink


class PPTXBuilder:
    """
    Simplified PowerPoint builder with one message per slide design.
    """

    def __init__(
        self,
        result: Any,
        charts: Dict[str, str],
        insights: Dict[str, str],
        language: str = "en"
    ):
        self.result = result
        self.charts = charts or {}
        self.insights = insights or {}
        self.language = language
        self.texts = get_all_texts(language)

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)

        self.is_comprehensive = hasattr(result, 'incident_summary')
        self.slide_count = 0

    # =========================================================================
    # UTILITY METHODS
    # =========================================================================

    def _set_background(self, slide) -> None:
        """Set dark background."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def _add_page_number(self, slide) -> None:
        """Add page number to bottom right."""
        self.slide_count += 1
        box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3)
        )
        p = box.text_frame.paragraphs[0]
        p.text = str(self.slide_count)
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.RIGHT

    def _add_title_bar(self, slide, title: str, accent_color: RGBColor = None) -> None:
        """Add minimal title bar at top."""
        # Thin accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color or PRIMARY
        line.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.3), Inches(12), Inches(0.5)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    def _add_insight_box(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        insight_text: str
    ) -> None:
        """Add AI insight box with professional styling."""
        # Box background
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = BG_ACCENT
        box.line.width = Pt(1)

        # AI label
        label_box = slide.shapes.add_textbox(
            Inches(left + 0.25), Inches(top + 0.2), Inches(2), Inches(0.3)
        )
        p = label_box.text_frame.paragraphs[0]
        p.text = "💡 AI INSIGHT" if self.language == "en" else "💡 AI 洞察"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        # Content
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.25), Inches(top + 0.55),
            Inches(width - 0.5), Inches(height - 0.75)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        # Smart truncation
        max_chars = int(width * 40)
        if len(insight_text) > max_chars:
            cut = insight_text.rfind(' ', 0, max_chars)
            insight_text = insight_text[:cut if cut > 0 else max_chars] + "..."
        p.text = insight_text
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_GRAY
        p.line_spacing = 1.3

    def _embed_chart(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        chart_key: str
    ) -> bool:
        """Embed a chart image. Returns True if successful."""
        if chart_key not in self.charts or not self.charts[chart_key]:
            return False
        try:
            img_data = base64.b64decode(self.charts[chart_key])
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(
                img_stream,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            return True
        except Exception:
            return False

    def _add_kpi_box(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        value: str,
        label: str,
        color: RGBColor = None
    ) -> None:
        """Add a KPI display box."""
        # Background
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.fill.background()

        # Top accent line
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.04)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color or PRIMARY
        accent.line.fill.background()

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + 0.25), Inches(width - 0.3), Inches(0.6)
        )
        p = value_box.text_frame.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(left + 0.15), Inches(top + height - 0.45), Inches(width - 0.3), Inches(0.3)
        )
        p = label_box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_GRAY

    # =========================================================================
    # SLIDE BUILDERS
    # =========================================================================

    def build_slide_cover(self) -> None:
        """Slide 1: Clean cover with title and date only."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # Top accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = PRIMARY
        line.line.fill.background()

        # Main title
        title = self.texts.get("report_title", "Operation Quality Report")
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.8), Inches(11.733), Inches(0.8)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.6), Inches(11.733), Inches(0.5)
        )
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = self.texts.get("report_subtitle", "Comprehensive ITIL Service Dashboard")
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_GRAY

        # Date range
        date_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.4), Inches(11.733), Inches(0.4)
        )
        p = date_box.text_frame.paragraphs[0]
        p.text = f"{self.result.start_date}  —  {self.result.end_date}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED

        # Bottom line
        bottom_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(5.2), Inches(4), Inches(0.02)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = BG_ACCENT
        bottom_line.line.fill.background()

        self._add_page_number(slide)

    def build_slide_overview(self) -> None:
        """Slide 2: Health score + 4 process overview (left-right layout)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = "Executive Summary" if self.language == "en" else "执行摘要"
        self._add_title_bar(slide, title)

        # LEFT SIDE: Health score + process cards
        # Health Score - Large
        score = self.result.health_score
        if score >= 90:
            score_color = SUCCESS
        elif score >= 80:
            score_color = WARNING
        else:
            score_color = DANGER

        score_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.1), Inches(5.8), Inches(2.2)
        )
        score_box.fill.solid()
        score_box.fill.fore_color.rgb = BG_CARD
        score_box.line.fill.background()

        # Score accent
        score_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.1), Inches(0.06), Inches(2.2)
        )
        score_accent.fill.solid()
        score_accent.fill.fore_color.rgb = score_color
        score_accent.line.fill.background()

        # Score label
        label_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.3), Inches(5), Inches(0.3)
        )
        p = label_box.text_frame.paragraphs[0]
        p.text = "HEALTH SCORE" if self.language == "en" else "健康评分"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_GRAY

        # Score value
        score_value_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.7), Inches(3), Inches(1)
        )
        p = score_value_box.text_frame.paragraphs[0]
        p.text = f"{score:.0f}"
        p.font.size = Pt(64)
        p.font.bold = True
        p.font.color.rgb = score_color

        # Grade
        grade_box = slide.shapes.add_textbox(
            Inches(3.2), Inches(2.1), Inches(3), Inches(0.5)
        )
        p = grade_box.text_frame.paragraphs[0]
        p.text = f"{self.result.health_emoji} {self.result.health_grade}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_WHITE

        # Process cards (2x2 grid)
        processes = [
            (self.texts.get("process_incident", "Incidents"), self.result.total_incidents, COLOR_INCIDENT),
            (self.texts.get("process_change", "Changes"), self.result.total_changes, COLOR_CHANGE),
            (self.texts.get("process_request", "Requests"), self.result.total_requests, COLOR_REQUEST),
            (self.texts.get("process_problem", "Problems"), self.result.total_problems, COLOR_PROBLEM),
        ]

        card_w, card_h = 2.8, 1.3
        positions = [
            (0.6, 3.5), (3.5, 3.5),
            (0.6, 4.9), (3.5, 4.9)
        ]

        for i, (name, count, color) in enumerate(processes):
            left, top = positions[i]
            self._add_kpi_box(slide, left, top, card_w, card_h, str(count), name, color)

        # RIGHT SIDE: AI Insight
        insight_text = self.insights.get(
            "executive_summary",
            self.texts.get("no_ai_insight", "AI insight not available")
        )
        self._add_insight_box(slide, 6.7, 1.1, 6.0, 5.1, insight_text)

        self._add_page_number(slide)

    def build_slide_incident(self) -> None:
        """Slide 3: Incident management (top-bottom layout)."""
        if not self.is_comprehensive:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = self.texts.get("process_incident", "Incident Management")
        self._add_title_bar(slide, title, COLOR_INCIDENT)

        # TOP: Chart area
        chart_embedded = self._embed_chart(slide, 0.6, 1.1, 12.133, 3.2, "sla_trend")

        if not chart_embedded:
            # Fallback: Show KPI cards
            kpis = self.result.incident_summary.kpis if self.result.incident_summary else {}
            kpi_data = [
                ("SLA Rate", kpis.get("sla_rate"), lambda v: f"{v.current_value:.1%}"),
                ("Avg MTTR", kpis.get("avg_mttr"), lambda v: f"{v.current_value:.1f}h"),
                ("Total", kpis.get("total_incidents"), lambda v: f"{int(v.current_value)}"),
                ("P1/P2", kpis.get("p1_p2_count"), lambda v: f"{int(v.current_value)}"),
            ]
            card_w = 2.9
            for i, (label, kpi, fmt) in enumerate(kpi_data):
                if kpi:
                    self._add_kpi_box(
                        slide, 0.6 + i * (card_w + 0.2), 1.1,
                        card_w, 1.4, fmt(kpi), label, COLOR_INCIDENT
                    )

        # BOTTOM: AI Insight
        insight_text = self.insights.get(
            "incident_insight",
            self.texts.get("no_ai_insight", "AI insight not available")
        )
        self._add_insight_box(slide, 0.6, 4.5, 12.133, 2.5, insight_text)

        self._add_page_number(slide)

    def build_slide_change(self) -> None:
        """Slide 4: Change management (top-bottom layout)."""
        if not self.is_comprehensive or not self.result.change_summary:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = self.texts.get("process_change", "Change Management")
        self._add_title_bar(slide, title, COLOR_CHANGE)

        # TOP: KPI cards (changes typically don't have trend charts)
        kpis = self.result.change_summary.kpis if self.result.change_summary else {}
        kpi_data = [
            ("Success Rate" if self.language == "en" else "成功率",
             kpis.get("change_success_rate"), lambda v: f"{v.current_value:.1%}"),
            ("Incident Rate" if self.language == "en" else "引发事件",
             kpis.get("change_incident_rate"), lambda v: f"{v.current_value:.1%}"),
            ("Total" if self.language == "en" else "总数",
             kpis.get("total_changes"), lambda v: f"{int(v.current_value)}"),
            ("Emergency" if self.language == "en" else "紧急变更",
             kpis.get("emergency_rate"), lambda v: f"{v.current_value:.1%}"),
        ]

        card_w = 2.9
        for i, (label, kpi, fmt) in enumerate(kpi_data):
            if kpi:
                status_color = COLOR_CHANGE
                if hasattr(kpi, 'status'):
                    if kpi.status == 'danger':
                        status_color = DANGER
                    elif kpi.status == 'warning':
                        status_color = WARNING
                self._add_kpi_box(
                    slide, 0.6 + i * (card_w + 0.2), 1.1,
                    card_w, 1.6, fmt(kpi), label, status_color
                )

        # Failed changes summary
        if self.result.failed_changes:
            failed_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.6), Inches(2.9), Inches(12.133), Inches(1.3)
            )
            failed_box.fill.solid()
            failed_box.fill.fore_color.rgb = BG_CARD
            failed_box.line.fill.background()

            failed_title = slide.shapes.add_textbox(
                Inches(0.8), Inches(3.0), Inches(4), Inches(0.3)
            )
            p = failed_title.text_frame.paragraphs[0]
            p.text = "Failed Changes" if self.language == "en" else "失败变更"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = DANGER

            failed_list = slide.shapes.add_textbox(
                Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.8)
            )
            tf = failed_list.text_frame
            tf.word_wrap = True
            items = [f"• {chg.change_number}: {chg.title[:40]}" for chg in self.result.failed_changes[:3]]
            p = tf.paragraphs[0]
            p.text = "  |  ".join(items) if items else "None"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_GRAY

        # BOTTOM: AI Insight
        insight_text = self.insights.get(
            "change_insight",
            self.texts.get("no_ai_insight", "AI insight not available")
        )
        self._add_insight_box(slide, 0.6, 4.5, 12.133, 2.5, insight_text)

        self._add_page_number(slide)

    def build_slide_request(self) -> None:
        """Slide 5: Service request (left-right layout)."""
        if not self.is_comprehensive or not self.result.request_summary:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = self.texts.get("process_request", "Service Requests")
        self._add_title_bar(slide, title, COLOR_REQUEST)

        # LEFT: KPIs
        kpis = self.result.request_summary.kpis if self.result.request_summary else {}

        # CSAT score prominent
        csat = kpis.get("request_csat")
        if csat:
            csat_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.6), Inches(1.1), Inches(5.8), Inches(2.0)
            )
            csat_box.fill.solid()
            csat_box.fill.fore_color.rgb = BG_CARD
            csat_box.line.fill.background()

            # CSAT value
            csat_val = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.3), Inches(3), Inches(0.8)
            )
            p = csat_val.text_frame.paragraphs[0]
            p.text = f"{csat.current_value:.1f}/5"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = COLOR_REQUEST

            # Stars visualization
            stars = "★" * int(csat.current_value) + "☆" * (5 - int(csat.current_value))
            stars_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.3), Inches(5), Inches(0.4)
            )
            p = stars_box.text_frame.paragraphs[0]
            p.text = stars + "  CSAT Score"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_GRAY

        # Other KPIs
        other_kpis = [
            ("Request SLA" if self.language == "en" else "请求 SLA",
             kpis.get("request_sla_rate"), lambda v: f"{v.current_value:.1%}"),
            ("Total" if self.language == "en" else "总数",
             kpis.get("total_requests"), lambda v: f"{int(v.current_value)}"),
        ]
        for i, (label, kpi, fmt) in enumerate(other_kpis):
            if kpi:
                self._add_kpi_box(
                    slide, 0.6 + i * 3.0, 3.3,
                    2.8, 1.2, fmt(kpi), label, COLOR_REQUEST
                )

        # RIGHT: AI Insight
        insight_text = self.insights.get(
            "request_insight",
            self.texts.get("no_ai_insight", "AI insight not available")
        )
        self._add_insight_box(slide, 6.7, 1.1, 6.0, 5.4, insight_text)

        self._add_page_number(slide)

    def build_slide_problem(self) -> None:
        """Slide 6: Problem management (top-bottom layout)."""
        if not self.is_comprehensive or not self.result.problem_summary:
            return

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = self.texts.get("process_problem", "Problem Management")
        self._add_title_bar(slide, title, COLOR_PROBLEM)

        # TOP: KPIs + Open problems count
        kpis = self.result.problem_summary.kpis if self.result.problem_summary else {}

        # Large open problems display
        open_count = len(self.result.open_problems) if self.result.open_problems else 0
        open_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.1), Inches(3.5), Inches(2.0)
        )
        open_box.fill.solid()
        open_box.fill.fore_color.rgb = BG_CARD
        open_box.line.fill.background()

        open_val = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3), Inches(3), Inches(0.8)
        )
        p = open_val.text_frame.paragraphs[0]
        p.text = str(open_count)
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = DANGER if open_count > 10 else COLOR_PROBLEM

        open_label = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.3), Inches(3), Inches(0.4)
        )
        p = open_label.text_frame.paragraphs[0]
        p.text = "Open Problems" if self.language == "en" else "待解决问题"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_GRAY

        # Other KPIs
        kpi_data = [
            ("Closure Rate" if self.language == "en" else "关闭率",
             kpis.get("problem_closure_rate"), lambda v: f"{v.current_value:.1%}"),
            ("RCA Rate" if self.language == "en" else "RCA完成",
             kpis.get("rca_rate"), lambda v: f"{v.current_value:.1%}"),
        ]
        for i, (label, kpi, fmt) in enumerate(kpi_data):
            if kpi:
                self._add_kpi_box(
                    slide, 4.3 + i * 3.0, 1.1,
                    2.8, 1.5, fmt(kpi), label, COLOR_PROBLEM
                )

        # Open problems list
        if self.result.open_problems:
            list_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.3), Inches(2.8), Inches(8.4), Inches(1.3)
            )
            list_box.fill.solid()
            list_box.fill.fore_color.rgb = BG_CARD
            list_box.line.fill.background()

            list_text = slide.shapes.add_textbox(
                Inches(4.5), Inches(2.9), Inches(8), Inches(1.1)
            )
            tf = list_text.text_frame
            tf.word_wrap = True
            items = []
            for prb in self.result.open_problems[:3]:
                title_short = prb.title[:35] + "..." if len(prb.title) > 35 else prb.title
                items.append(f"• {prb.problem_number}: {title_short}")
            p = tf.paragraphs[0]
            p.text = "\n".join(items)
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_GRAY
            p.line_spacing = 1.4

        # BOTTOM: AI Insight
        insight_text = self.insights.get(
            "problem_insight",
            self.texts.get("no_ai_insight", "AI insight not available")
        )
        self._add_insight_box(slide, 0.6, 4.4, 12.133, 2.6, insight_text)

        self._add_page_number(slide)

    def build_slide_trends(self) -> None:
        """Slide 7: Trend analysis (top-bottom layout)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = "Trend Analysis" if self.language == "en" else "趋势分析"
        self._add_title_bar(slide, title)

        # TOP: Charts (2 side by side)
        left_embedded = self._embed_chart(slide, 0.6, 1.1, 6.0, 3.0, "incident_volume")
        right_embedded = self._embed_chart(slide, 6.8, 1.1, 6.0, 3.0, "health_gauge")

        if not left_embedded and not right_embedded:
            # Fallback message
            msg_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(2.0), Inches(12), Inches(1)
            )
            p = msg_box.text_frame.paragraphs[0]
            p.text = "Charts require visualization data" if self.language == "en" else "图表需要可视化数据"
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MUTED
            p.alignment = PP_ALIGN.CENTER

        # BOTTOM: AI Insight
        insight_text = self.insights.get(
            "trend_insight",
            self.insights.get("executive_summary", self.texts.get("no_ai_insight", ""))
        )
        self._add_insight_box(slide, 0.6, 4.3, 12.133, 2.7, insight_text)

        self._add_page_number(slide)

    def build_slide_risks(self) -> None:
        """Slide 8: Risk radar (left-right layout)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = "Risk Radar" if self.language == "en" else "风险雷达"
        self._add_title_bar(slide, title, DANGER)

        # LEFT: Risk list
        risks = self.result.top_risks[:4] if self.result.top_risks else []

        if risks:
            for i, risk in enumerate(risks):
                top = 1.2 + i * 1.4
                
                # Risk card
                risk_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.6), Inches(top), Inches(5.8), Inches(1.25)
                )
                risk_bg.fill.solid()
                risk_bg.fill.fore_color.rgb = BG_CARD
                risk_bg.line.fill.background()

                # Priority indicator
                priority_color = DANGER if risk.priority == "Critical" else (WARNING if risk.priority == "Warning" else TEXT_MUTED)
                indicator = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.6), Inches(top), Inches(0.05), Inches(1.25)
                )
                indicator.fill.solid()
                indicator.fill.fore_color.rgb = priority_color
                indicator.line.fill.background()

                # Priority badge
                badge = slide.shapes.add_textbox(
                    Inches(0.8), Inches(top + 0.1), Inches(1.5), Inches(0.25)
                )
                p = badge.text_frame.paragraphs[0]
                p.text = risk.priority.upper()
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = priority_color

                # Message
                msg = slide.shapes.add_textbox(
                    Inches(0.8), Inches(top + 0.4), Inches(5.4), Inches(0.8)
                )
                tf = msg.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = risk.message[:80] if len(risk.message) > 80 else risk.message
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_WHITE
        else:
            no_risk = slide.shapes.add_textbox(
                Inches(0.6), Inches(3), Inches(5.8), Inches(1)
            )
            p = no_risk.text_frame.paragraphs[0]
            p.text = "✓ No significant risks identified" if self.language == "en" else "✓ 未发现重大风险"
            p.font.size = Pt(14)
            p.font.color.rgb = SUCCESS

        # RIGHT: AI Insight
        insight_text = self.insights.get(
            "risk_insight",
            self.insights.get("executive_summary", self.texts.get("no_ai_insight", ""))
        )
        self._add_insight_box(slide, 6.7, 1.2, 6.0, 5.2, insight_text)

        self._add_page_number(slide)

    def build_slide_actions(self) -> None:
        """Slide 9: Recommended actions (left-right layout)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        title = "Recommended Actions" if self.language == "en" else "改进建议"
        self._add_title_bar(slide, title, SUCCESS)

        # LEFT: Action list
        actions = self.result.actions[:4] if self.result.actions else []

        if actions:
            for i, action in enumerate(actions):
                top = 1.2 + i * 1.4

                # Action card
                action_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.6), Inches(top), Inches(5.8), Inches(1.25)
                )
                action_bg.fill.solid()
                action_bg.fill.fore_color.rgb = BG_CARD
                action_bg.line.fill.background()

                # Priority indicator
                priority_color = DANGER if action.priority == "Urgent" else (WARNING if action.priority == "High" else SUCCESS)
                indicator = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.6), Inches(top), Inches(0.05), Inches(1.25)
                )
                indicator.fill.solid()
                indicator.fill.fore_color.rgb = priority_color
                indicator.line.fill.background()

                # Priority badge
                badge = slide.shapes.add_textbox(
                    Inches(0.8), Inches(top + 0.1), Inches(1.5), Inches(0.25)
                )
                p = badge.text_frame.paragraphs[0]
                p.text = action.priority.upper()
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = priority_color

                # Action text
                action_text = slide.shapes.add_textbox(
                    Inches(0.8), Inches(top + 0.4), Inches(5.4), Inches(0.5)
                )
                tf = action_text.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = action.action[:70] if len(action.action) > 70 else action.action
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_WHITE

                # Expected impact
                impact = slide.shapes.add_textbox(
                    Inches(0.8), Inches(top + 0.9), Inches(5.4), Inches(0.3)
                )
                p = impact.text_frame.paragraphs[0]
                p.text = f"↗ {action.expected_impact[:50]}"
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT_MUTED
        else:
            no_actions = slide.shapes.add_textbox(
                Inches(0.6), Inches(3), Inches(5.8), Inches(1)
            )
            p = no_actions.text_frame.paragraphs[0]
            p.text = "No actions required at this time" if self.language == "en" else "目前无需采取行动"
            p.font.size = Pt(14)
            p.font.color.rgb = SUCCESS

        # RIGHT: AI Insight
        insight_text = self.insights.get(
            "action_insight",
            self.insights.get("executive_summary", self.texts.get("no_ai_insight", ""))
        )
        self._add_insight_box(slide, 6.7, 1.2, 6.0, 5.2, insight_text)

        self._add_page_number(slide)

    def build_slide_closing(self) -> None:
        """Slide 10: Clean closing slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)

        # Center line
        center_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(3.5), Inches(13.333), Inches(0.5)
        )
        center_line.fill.solid()
        center_line.fill.fore_color.rgb = PRIMARY
        center_line.line.fill.background()

        # Thank you
        thanks = slide.shapes.add_textbox(
            Inches(0), Inches(2.5), Inches(13.333), Inches(0.8)
        )
        p = thanks.text_frame.paragraphs[0]
        p.text = "Thank You" if self.language == "en" else "谢谢"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

        # Generated time
        time_box = slide.shapes.add_textbox(
            Inches(0), Inches(4.2), Inches(13.333), Inches(0.4)
        )
        p = time_box.text_frame.paragraphs[0]
        gen_text = f"{self.texts.get('generated_on', 'Generated')}: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        p.text = gen_text
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER

        # Confidential
        conf_box = slide.shapes.add_textbox(
            Inches(0), Inches(6.8), Inches(13.333), Inches(0.4)
        )
        p = conf_box.text_frame.paragraphs[0]
        p.text = self.texts.get("footer_confidential", "Confidential - Internal Use Only")
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER

        self._add_page_number(slide)

    # =========================================================================
    # BUILD & SAVE
    # =========================================================================

    def build(self) -> None:
        """Build complete presentation."""
        self.build_slide_cover()
        self.build_slide_overview()

        if self.is_comprehensive:
            self.build_slide_incident()
            self.build_slide_change()
            self.build_slide_request()
            self.build_slide_problem()

        self.build_slide_trends()
        self.build_slide_risks()
        self.build_slide_actions()
        self.build_slide_closing()

    def save(self, filepath: Path = None) -> str:
        """Save presentation to file."""
        if filepath is None:
            lang = "CN" if self.language == "zh" else "EN"
            filename = f"Operation_Quality_Report_{self.result.start_date}_to_{self.result.end_date}_{lang}.pptx"
            filepath = OUTPUT_DIR / filename

        self.build()
        self.prs.save(filepath)
        return str(filepath)
